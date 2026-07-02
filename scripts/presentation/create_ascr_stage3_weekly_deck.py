#!/usr/bin/env python3
"""Generate a concise bilingual Stage-3 weekly report deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "presentations" / "ascr_stage3_weekly_report.pptx"

NAVY = RGBColor(32, 48, 68)
SLATE = RGBColor(70, 82, 96)
MUTED = RGBColor(105, 115, 128)
LIGHT_BG = RGBColor(247, 248, 250)
LINE = RGBColor(205, 211, 219)
BLUE = RGBColor(66, 103, 146)
TEAL = RGBColor(78, 129, 124)
AMBER = RGBColor(156, 122, 63)
GREEN = RGBColor(88, 135, 96)
RED = RGBColor(158, 91, 91)
WHITE = RGBColor(255, 255, 255)


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def set_text(frame, text: str, size: int = 18, color=SLATE, bold: bool = False):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.6))
    set_text(box.text_frame, title, 24, NAVY, True)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.88), Inches(12.1), Inches(0.32))
        set_text(sub.text_frame, subtitle, 10, MUTED, False)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page: int):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(11.8), Inches(0.25))
    set_text(box.text_frame, f"ASCR Stage 3 weekly update | Page {page}", 8, MUTED)


def bullet_box(slide, x, y, w, h, title: str, bullets: list[str], accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.8)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.18), w - Inches(0.42), h - Inches(0.25))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NAVY
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = SLATE
        p.space_before = Pt(6)


def flow_node(slide, x, y, w, h, text: str, fill=LIGHT_BG, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.9)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(shape.text_frame, text, 12, color, True)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def arrow(slide, x1, y1, x2, y2, color=SLATE):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(1.6)
    conn.line.end_arrowhead = True
    return conn


def label(slide, x, y, w, h, text: str, size=10, color=MUTED, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(x, y, w, h)
    set_text(box.text_frame, text, size, color)
    box.text_frame.paragraphs[0].alignment = align
    return box


def phase_chip(slide, x, y, name: str, status: str, color):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.65), Inches(0.48))
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(chip.text_frame, f"{name}\n{status}", 9, WHITE, True)
    chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb("F7F8FA")
    return slide


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = add_slide(prs)
    label(slide, Inches(0.7), Inches(0.55), Inches(3.0), Inches(0.25), "PhD weekly report / 组会汇报", 10, MUTED, PP_ALIGN.LEFT)
    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(11.5), Inches(1.0))
    set_text(title.text_frame, "Stage 3: Self-Corrupted Token Repair", 34, NAVY, True)
    subtitle = slide.shapes.add_textbox(Inches(0.72), Inches(2.35), Inches(10.6), Inches(0.6))
    set_text(subtitle.text_frame, "Discrete-token corruption → clean-token/image recovery → SFT pairs → Lumina repair behavior", 17, SLATE)
    bullet_box(
        slide,
        Inches(0.72),
        Inches(3.25),
        Inches(5.7),
        Inches(1.8),
        "本周汇报范围",
        ["Only Stage 3, Phase 0-5", "不展开 Stage 1/2 细节", "不编造未记录实验结果"],
        BLUE,
    )
    bullet_box(
        slide,
        Inches(6.75),
        Inches(3.25),
        Inches(5.5),
        Inches(1.8),
        "Core claim / 核心问题",
        ["Can corrupted discrete VQ tokens teach localized repair?", "Labels come from known token masks, not human judgment"],
        TEAL,
    )
    add_footer(slide, 1)

    # 2
    slide = add_slide(prs)
    add_title(slide, "Stage 3 Objective / 目标", "From controlled token corruption to supervised repair behavior")
    nodes = [
        ("Clean 64×64\nVQ tokens", BLUE),
        ("Controlled\ncorruption", AMBER),
        ("Corrupted\nVQ tokens", RED),
        ("Known mask →\n8×8 cells", TEAL),
        ("SFT target\n{\"cells\": [...]}", GREEN),
    ]
    x = Inches(0.75)
    y = Inches(2.2)
    for i, (text, color) in enumerate(nodes):
        flow_node(slide, x + Inches(i * 2.45), y, Inches(1.85), Inches(0.82), text, WHITE, color)
        if i < len(nodes) - 1:
            arrow(slide, x + Inches(i * 2.45 + 1.85), y + Inches(0.41), x + Inches(i * 2.45 + 2.35), y + Inches(0.41))
    bullet_box(
        slide,
        Inches(0.85),
        Inches(4.1),
        Inches(11.6),
        Inches(1.45),
        "What Stage 3 is trying to learn",
        [
            "Input: original prompt + corrupted discrete image-token state",
            "Output: compact repair-cell JSON, e.g. {\"cells\":[\"D4\",\"D5\"]}",
            "Training signal: corrupted-clean pairs generated inside Lumina token space",
        ],
        BLUE,
    )
    add_footer(slide, 2)

    # 3
    slide = add_slide(prs)
    add_title(slide, "Phase Map / 阶段地图", "Phase 0-5 only")
    xs = [0.8, 2.72, 4.64, 6.56, 8.48, 10.4]
    phases = [
        ("P0", "Direction reset", BLUE),
        ("P1", "Locality probe", GREEN),
        ("P2", "Dataset build", GREEN),
        ("P3", "Selector baselines", GREEN),
        ("P4", "MMU/LoRA", AMBER),
        ("P5", "ASCR loop", SLATE),
    ]
    for i, (p, text, color) in enumerate(phases):
        phase_chip(slide, Inches(xs[i]), Inches(1.75), p, text, color)
        if i < len(phases) - 1:
            arrow(slide, Inches(xs[i] + 1.65), Inches(1.99), Inches(xs[i + 1] - 0.08), Inches(1.99), LINE)
    bullet_box(slide, Inches(0.85), Inches(3.0), Inches(3.75), Inches(2.4), "Completed / 已完成", ["P0 research direction reset", "P1 smoke locality gate", "P2 self-corruption dataset tooling", "P3 baseline gate on Hard64"], GREEN)
    bullet_box(slide, Inches(4.8), Inches(3.0), Inches(3.75), Inches(2.4), "In progress / 进行中", ["P4 repair_cells SFT + LoRA pipeline", "1024px H200/L40S training profiles", "Parse-rate-first evaluation"], AMBER)
    bullet_box(slide, Inches(8.75), Inches(3.0), Inches(3.75), Inches(2.4), "Next / 下一步", ["P5 loop integration after useful selector", "Synthetic localization + image-quality evaluation", "No claims beyond current evidence"], BLUE)
    add_footer(slide, 3)

    # 4
    slide = add_slide(prs)
    add_title(slide, "Phase 0: Direction Reset / 方向重置", "Old cross-model Stage 3 is no longer the mainline")
    bullet_box(slide, Inches(0.8), Inches(1.65), Inches(5.6), Inches(2.0), "Research framing", ["Mainline: self-corrupted token repair", "Conservative claim: local repair signal, not general intelligence", "Stage 2 JSON evaluator remains separate"], BLUE)
    bullet_box(slide, Inches(6.9), Inches(1.65), Inches(5.6), Inches(2.0), "Repository artifacts", ["Design doc updated", "Server handoff docs added", "Shared collaboration log records decisions"], TEAL)
    flow_node(slide, Inches(1.1), Inches(4.55), Inches(2.2), Inches(0.75), "Prompt + image\nStage 2 path", WHITE, MUTED)
    flow_node(slide, Inches(5.55), Inches(4.55), Inches(2.2), Inches(0.75), "Prompt + corrupted\nVQ tokens", WHITE, BLUE)
    flow_node(slide, Inches(10.0), Inches(4.55), Inches(2.2), Inches(0.75), "TokenReopenMask\nStage 3 path", WHITE, TEAL)
    arrow(slide, Inches(3.35), Inches(4.93), Inches(5.45), Inches(4.93))
    arrow(slide, Inches(7.8), Inches(4.93), Inches(9.9), Inches(4.93))
    add_footer(slide, 4)

    # 5
    slide = add_slide(prs)
    add_title(slide, "Phase 1: Token Locality Probe / 局部性验证", "Before training, check whether token corruption has local visual effects")
    labels = ["Clean tokens", "Corrupt block", "Decode", "Difference map", "Locality metrics"]
    for i, text in enumerate(labels):
        flow_node(slide, Inches(0.85 + i * 2.4), Inches(1.75), Inches(1.85), Inches(0.72), text, WHITE, [BLUE, AMBER, TEAL, RED, GREEN][i])
        if i < len(labels) - 1:
            arrow(slide, Inches(2.7 + i * 2.4), Inches(2.11), Inches(3.15 + i * 2.4), Inches(2.11))
    bullet_box(slide, Inches(0.85), Inches(3.35), Inches(5.6), Inches(2.3), "Implemented", ["Corruptors: random_replace, local_shuffle, neighbor_copy, transplant", "Metrics: inside energy, top-1/top-k hit, radius", "Outputs: manifest, summary, heatmaps, images, tokens"], BLUE)
    bullet_box(slide, Inches(6.9), Inches(3.35), Inches(5.6), Inches(2.3), "Recorded result", ["Smoke: 8 prompts, 24 corruption rows", "Generation/decode succeeded for all rows", "4×4 random replace and local_shuffle showed clear locality"], GREEN)
    add_footer(slide, 5)

    # 6
    slide = add_slide(prs)
    add_title(slide, "Phase 2: Self-Corruption Dataset / 数据构造", "Create corrupted-clean pairs with exact labels")
    for i, (t, c) in enumerate([
        ("clean_vq_ids\n64×64", BLUE),
        ("mask\n1/2/4/8", AMBER),
        ("operator\n4 choices", AMBER),
        ("corrupted_vq_ids", RED),
        ("project to\n8×8 cells", TEAL),
        ("SFT row\nJSON target", GREEN),
    ]):
        flow_node(slide, Inches(0.55 + i * 2.08), Inches(1.6), Inches(1.58), Inches(0.68), t, WHITE, c)
        if i < 5:
            arrow(slide, Inches(2.13 + i * 2.08), Inches(1.94), Inches(2.55 + i * 2.08), Inches(1.94))
    bullet_box(slide, Inches(0.8), Inches(3.0), Inches(5.55), Inches(2.55), "Canonical dataset row", ["Positive: corrupted tokens + nonempty cells", "Negative: same clean tokens + {\"cells\":[]}", "Decoded images are audit-only, not label sources"], BLUE)
    bullet_box(slide, Inches(6.85), Inches(3.0), Inches(5.55), Inches(2.55), "Current scale target", ["10k clean prompts / clean tokens", "30k positive + 10k negative repair rows", "Group-safe SFT split by source clean sample"], TEAL)
    add_footer(slide, 6)

    # 7
    slide = add_slide(prs)
    add_title(slide, "Phase 3: Selector Baselines / 基线选择器", "Check learnability before training Lumina adapters")
    bullet_box(slide, Inches(0.8), Inches(1.55), Inches(3.7), Inches(2.6), "Baselines", ["random", "token_prior", "RGB diff oracle", "RGB localizer", "prompt + RGB localizer"], BLUE)
    bullet_box(slide, Inches(4.85), Inches(1.55), Inches(3.7), Inches(2.6), "Metrics", ["precision / recall / F1", "cell IoU", "hit_any", "distance to target cells"], TEAL)
    bullet_box(slide, Inches(8.9), Inches(1.55), Inches(3.7), Inches(2.6), "Recorded gate", ["Hard64 dataset: 128 rows", "prompt_rgb_localizer: 0.875 hit_any at 16×16", "Enough to justify Phase 4"], GREEN)
    label(slide, Inches(1.0), Inches(5.05), Inches(11.3), Inches(0.6), "Important: the 24-row smoke set validates wiring only; stronger conclusions rely on expanded datasets.", 13, MUTED)
    add_footer(slide, 7)

    # 8
    slide = add_slide(prs)
    add_title(slide, "Phase 4: MMU / LoRA Repair Learning", "Train Lumina to output repair cells from corrupted token state")
    flow_node(slide, Inches(0.85), Inches(1.65), Inches(2.0), Inches(0.72), "Prompt +\ncorrupted tokens", WHITE, BLUE)
    flow_node(slide, Inches(3.55), Inches(1.65), Inches(2.0), Inches(0.72), "Lumina MMU\nanswer path", WHITE, TEAL)
    flow_node(slide, Inches(6.25), Inches(1.65), Inches(2.0), Inches(0.72), "LoRA SFT\nadapter", WHITE, AMBER)
    flow_node(slide, Inches(8.95), Inches(1.65), Inches(2.0), Inches(0.72), "JSON\n{\"cells\":[]}", WHITE, GREEN)
    arrow(slide, Inches(2.85), Inches(2.01), Inches(3.45), Inches(2.01))
    arrow(slide, Inches(5.55), Inches(2.01), Inches(6.15), Inches(2.01))
    arrow(slide, Inches(8.25), Inches(2.01), Inches(8.85), Inches(2.01))
    bullet_box(slide, Inches(0.8), Inches(3.05), Inches(5.55), Inches(2.45), "What has been implemented", ["SFT prep + Lumina SFT conversion", "VQ-token and decoded-image input modes", "Single-GPU and DDP LoRA training wrappers", "H200/L40S profile configs"], BLUE)
    bullet_box(slide, Inches(6.85), Inches(3.05), Inches(5.55), Inches(2.45), "Known Phase-4 evidence", ["Zero-shot parse_rate: 0.0 on early probe", "Old-schema LoRA loss: 9.75 → 0.157", "Old-schema LoRA eval parse_rate: 0.156; hit_any: 0.0"], AMBER)
    add_footer(slide, 8)

    # 9
    slide = add_slide(prs)
    add_title(slide, "Current Mainline Update / 当前主线", "Token-only 8×8 repair_cells schema")
    bullet_box(slide, Inches(0.8), Inches(1.55), Inches(5.6), Inches(2.55), "Schema simplification", ["Target JSON has exactly one key: cells", "Positive example: {\"cells\":[\"D4\",\"D5\"]}", "Negative example: {\"cells\":[]}"], GREEN)
    bullet_box(slide, Inches(6.9), Inches(1.55), Inches(5.6), Inches(2.55), "Why this matters", ["Removes legacy key mismatch", "Directly matches reopen decision", "Parse rate must be near 1 before localization metrics matter"], BLUE)
    bullet_box(slide, Inches(0.8), Inches(4.55), Inches(11.7), Inches(1.05), "Resource note", ["1024px LoRA training is memory-sensitive; current configs include H200 1024px and L40S gradient-checkpointing/8-bit-Adam profiles."], AMBER)
    add_footer(slide, 9)

    # 10
    slide = add_slide(prs)
    add_title(slide, "Phase 5: ASCR Loop Integration / 闭环修复", "Only after a useful selector or LoRA localizer")
    nodes = [
        ("Lumina\nbaseline image", BLUE),
        ("Predict\nrepair cells", TEAL),
        ("Reopen selected\nVQ tokens", AMBER),
        ("Decode repaired\nimage", GREEN),
        ("Optional\naccept/reject", SLATE),
    ]
    for i, (t, c) in enumerate(nodes):
        flow_node(slide, Inches(0.85 + i * 2.4), Inches(1.75), Inches(1.85), Inches(0.72), t, WHITE, c)
        if i < len(nodes) - 1:
            arrow(slide, Inches(2.7 + i * 2.4), Inches(2.11), Inches(3.15 + i * 2.4), Inches(2.11))
    bullet_box(slide, Inches(0.85), Inches(3.45), Inches(5.6), Inches(2.0), "Initial policy", ["If mask empty: stop", "If mask nonempty: reopen selected cells", "Add scorer only after nontrivial masks"], BLUE)
    bullet_box(slide, Inches(6.9), Inches(3.45), Inches(5.6), Inches(2.0), "Artifacts to keep", ["Trace and selected indices", "Selected token/cell count", "Decoded before/after images", "Final summary metrics"], TEAL)
    add_footer(slide, 10)

    # 11
    slide = add_slide(prs)
    add_title(slide, "This Week's Deliverables / 本周完成", "Evidence-backed summary")
    bullet_box(slide, Inches(0.8), Inches(1.55), Inches(3.7), Inches(3.4), "Research design", ["Stage 3 scoped to discrete-token corruption", "Phase 0-5 roadmap clarified", "Training labels defined from token masks"], BLUE)
    bullet_box(slide, Inches(4.85), Inches(1.55), Inches(3.7), Inches(3.4), "Pipeline", ["Clean-token generation path", "Repair dataset config: 40k rows target", "SFT conversion and LoRA runners"], TEAL)
    bullet_box(slide, Inches(8.9), Inches(1.55), Inches(3.7), Inches(3.4), "Validation gates", ["Locality smoke passed", "Hard64 selector baseline gate cleared", "Phase 4 not yet fully cleared"], AMBER)
    label(slide, Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.55), "No formal before/after image-quality result is claimed here unless already recorded in the repository.", 13, MUTED)
    add_footer(slide, 11)

    # 12
    slide = add_slide(prs)
    add_title(slide, "Next Actions / 下一步", "Concrete server continuation")
    bullet_box(slide, Inches(0.8), Inches(1.45), Inches(5.6), Inches(2.2), "Short-term", ["Generate/verify clean tokens", "Build repair_cells_40k dataset", "Prepare SFT and convert Lumina JSONL"], BLUE)
    bullet_box(slide, Inches(6.9), Inches(1.45), Inches(5.6), Inches(2.2), "Training gate", ["Train one authoritative LoRA adapter", "Probe parse_rate first", "Only interpret hit_any / IoU after parse is reliable"], TEAL)
    bullet_box(slide, Inches(0.8), Inches(4.15), Inches(11.7), Inches(1.45), "Open risk", ["If 8-GPU 1024px LoRA OOMs, inspect DDP log and use H200/L40S fallback profiles; do not train many independent adapters and merge them."], AMBER)
    add_footer(slide, 12)

    # 13
    slide = add_slide(prs)
    add_title(slide, "Evidence Sources / 依据文件", "Repository files used to avoid unsupported claims")
    bullet_box(slide, Inches(0.8), Inches(1.45), Inches(11.7), Inches(4.25), "Checked files", [
        "docs/STAGE3_SELF_CORRUPTED_TOKEN_REPAIR.md",
        "docs/SERVER_AI_TASK_STAGE3_TOKEN_REPAIR_DATASET_AND_LORA.md",
        "docs/SERVER_AI_TASK_STAGE3_SELF_CORRUPT_LOCALITY.md",
        "docs/SERVER_AI_TASK_STAGE3_SELF_CORRUPT_DATASET.md",
        "docs/SERVER_AI_TASK_STAGE3_SELF_CORRUPT_SELECTORS.md",
        "configs/stage3/self_corrupt/token_repair_40k.yaml",
        "configs/stage4/self_corrupt/mmu_sft_token_repair_8x8.yaml",
    ], SLATE)
    add_footer(slide, 13)

    return prs


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
